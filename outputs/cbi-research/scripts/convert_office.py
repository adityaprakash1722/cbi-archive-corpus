#!/usr/bin/env python3
"""Convert the DOCX / DOC / PPTX / ZIP files the PDF pipeline skipped.

Background
----------
``convert_pdfs.py`` filters ``format == "PDF"``. Everything else that the
crawler downloaded was therefore never converted, never indexed and never
searchable: 201 DOCX, 71 DOC, 49 ZIP and 2 PPTX, 387 MB in total.

Honesty about page anchors
--------------------------
The PDF corpus has true source-page anchors. These formats mostly do not have
pages, so every output records ``page_basis``:

  slide                 one pseudo-page per slide (a real structural unit)
  explicit-page-break   the author inserted hard page breaks; those are used
  archive-member        one pseudo-page per file inside a ZIP
  single-pseudo-page    no page structure exists; the whole document is page 1

A ``single-pseudo-page`` anchor means "this document", not "this page", and
must not be cited as a page reference.
"""

from __future__ import annotations

import argparse, csv, functools, hashlib, io, json, os, re, shutil, subprocess, sys, tempfile, time, zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

PIPELINE_VERSION = "office-0.3.0"
TEXTUAL_SUFFIXES = {".xsd", ".xml", ".txt", ".csv", ".json", ".md", ".htm", ".html", ".xsl", ".dtd", ".ini", ".cfg"}
MAX_MEMBER_CHARS = 40000
W_NS = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def atomic_text_write(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + f".{os.getpid()}.tmp")
    tmp.write_text(text, encoding="utf-8", newline="\n")
    os.replace(tmp, path)


def yaml_scalar(v) -> str:
    return json.dumps(v, ensure_ascii=False)


# --------------------------------------------------------------------- DOCX
def _word_text(element: ET.Element) -> str:
    """Read each OOXML text node once, preserving inline tabs and breaks."""
    pieces: list[str] = []
    for node in element.iter():
        if node.tag == W_NS + "t":
            pieces.append(node.text or "")
        elif node.tag == W_NS + "tab":
            pieces.append("\t")
        elif node.tag == W_NS + "br":
            pieces.append("\f" if node.get(W_NS + "type") == "page" else "\n")
        elif node.tag == W_NS + "lastRenderedPageBreak":
            pieces.append("\f")
    return "".join(pieces)


def _paragraph_markdown(paragraph: ET.Element) -> list[str]:
    text = _word_text(paragraph)
    style_node = paragraph.find(f"{W_NS}pPr/{W_NS}pStyle")
    style = (style_node.get(W_NS + "val") if style_node is not None else "") or ""
    page_before = paragraph.find(f"{W_NS}pPr/{W_NS}pageBreakBefore") is not None
    parts = text.split("\f")
    rendered: list[str] = []
    for part in parts:
        clean = part.strip()
        if clean:
            heading = re.search(r"heading\s*([1-9])", style, re.I)
            if heading:
                clean = "#" * min(int(heading.group(1)), 6) + " " + clean
            elif style.casefold() in {"title", "subtitle"}:
                clean = "# " + clean
        rendered.append(clean)
    if page_before:
        rendered.insert(0, "")
    return rendered


def _table_markdown(table: ET.Element) -> str:
    """Render direct XML cells once, avoiding python-docx merged-cell aliases."""
    rendered: list[str] = []
    for row in table.findall(W_NS + "tr"):
        cells: list[str] = []
        for cell in row.findall(W_NS + "tc"):
            merge = cell.find(f"{W_NS}tcPr/{W_NS}vMerge")
            continuing_merge = (
                merge is not None
                and (merge.get(W_NS + "val") or "continue").casefold() != "restart"
            )
            text = "" if continuing_merge else " ".join(_word_text(cell).replace("\f", " ").split())
            cells.append(text.replace("|", "\\|"))
            span = cell.find(f"{W_NS}tcPr/{W_NS}gridSpan")
            try:
                cells.extend([""] * max(int(span.get(W_NS + "val")) - 1, 0) if span is not None else [])
            except (TypeError, ValueError):
                pass
        if cells:
            rendered.append("| " + " | ".join(cells) + " |")
    return "\n".join(rendered)


def _append_block(pages: list[list[str]], block_parts: list[str]) -> None:
    """Append a block whose list boundaries represent explicit page breaks."""
    for index, part in enumerate(block_parts):
        if index:
            pages.append([])
        if part.strip():
            pages[-1].append(part.strip())


def _ancillary_part(zf: zipfile.ZipFile, name: str) -> tuple[str, int]:
    try:
        root = ET.fromstring(zf.read(name))
    except KeyError:
        return "", 0
    text = "\n".join(
        " ".join(_word_text(child).replace("\f", " ").split())
        for child in root
        if _word_text(child).strip()
    ).strip()
    visible = sum(len(node.text or "") for node in root.iter(W_NS + "t"))
    return text, visible


def convert_docx(path: Path) -> tuple[list[str], str, str, dict]:
    """Extract a DOCX in body order without duplicating merged table cells.

    The v5.1 extractor enumerated ``document.paragraphs`` and then
    ``document.tables``. That reordered every interleaved table and, because
    python-docx exposes a merged cell through several grid positions, expanded
    two forms about forty-fold. Walking the direct OOXML body children preserves
    order and gives every physical ``w:tc`` element exactly one rendering.
    """
    with zipfile.ZipFile(path) as zf:
        document_xml = zf.read("word/document.xml")
        root = ET.fromstring(document_xml)
        body = root.find(W_NS + "body")
        if body is None:
            raise ValueError("word/document.xml has no w:body")
        pages: list[list[str]] = [[]]
        for child in body:
            if child.tag == W_NS + "p":
                _append_block(pages, _paragraph_markdown(child))
            elif child.tag == W_NS + "tbl":
                table = _table_markdown(child)
                if table:
                    pages[-1].append(table)

        source_visible = sum(len(node.text or "") for node in root.iter(W_NS + "t"))
        ancillary_sections: list[str] = []
        seen_ancillary: set[str] = set()
        part_groups = (
            ("Headers", sorted(name for name in zf.namelist() if re.fullmatch(r"word/header\d+\.xml", name))),
            ("Footers", sorted(name for name in zf.namelist() if re.fullmatch(r"word/footer\d+\.xml", name))),
            ("Footnotes", ["word/footnotes.xml"]),
            ("Endnotes", ["word/endnotes.xml"]),
            ("Comments", ["word/comments.xml"]),
        )
        for label, names in part_groups:
            texts: list[str] = []
            for name in names:
                text, visible = _ancillary_part(zf, name)
                source_visible += visible
                if text and text not in seen_ancillary:
                    seen_ancillary.add(text)
                    texts.append(text)
            if texts:
                ancillary_sections.append(f"## {label}\n\n" + "\n\n".join(texts))
        if ancillary_sections:
            pages[-1].extend(ancillary_sections)

    text_pages = ["\n\n".join(page).strip() for page in pages]
    text_pages = [page for page in text_pages if page] or [""]
    output_characters = sum(len(page) for page in text_pages)
    diagnostics = {
        "source_text_characters": source_visible,
        "output_expansion_ratio": round(output_characters / max(source_visible, 1), 4),
    }
    basis = "explicit-page-break" if len(text_pages) > 1 else "single-pseudo-page"
    return text_pages, "ooxml-body-order 1.0", basis, diagnostics


# --------------------------------------------------------------------- PPTX
def convert_pptx(path: Path) -> tuple[list[str], str, str]:
    from pptx import Presentation
    presentation = Presentation(str(path))
    pages = []
    for number, slide in enumerate(presentation.slides, 1):
        lines = [f"# Slide {number}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append("**Speaker notes:** " + slide.notes_slide.notes_text_frame.text.strip())
        pages.append("\n\n".join(lines))
    return (pages or [""]), "python-pptx", "slide"


# ---------------------------------------------------------------- legacy DOC
@functools.lru_cache(maxsize=1)
def libreoffice_engine() -> tuple[str, str]:
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError("LibreOffice 'soffice' executable is required but was not found on PATH")
    result = subprocess.run([soffice, "--version"], capture_output=True, text=True, timeout=30)
    version = (result.stdout or result.stderr).strip().replace("LibreOffice ", "")
    if result.returncode or not version:
        raise RuntimeError("could not determine LibreOffice version")
    return soffice, "libreoffice-headless " + version


def convert_doc(path: Path) -> tuple[list[str], str, str]:
    """LibreOffice headless. Replaces the original longest-run-only OLE scraper,
    which kept a single text fragment and discarded the rest of the document."""
    with tempfile.TemporaryDirectory() as workdir:
        staged = Path(workdir) / (re.sub(r"[^A-Za-z0-9._-]", "_", path.name) or "doc.doc")
        staged.write_bytes(path.read_bytes())
        profile = Path(workdir) / "profile"
        soffice, engine = libreoffice_engine()
        result = subprocess.run(
            [soffice, "--headless", "--norestore", f"-env:UserInstallation=file://{profile}",
             "--convert-to", "txt:Text (encoded):UTF8", "--outdir", workdir, str(staged)],
            capture_output=True, timeout=120,
        )
        produced = list(Path(workdir).glob("*.txt"))
        if not produced:
            raise RuntimeError(f"soffice produced no output: {result.stderr.decode('utf-8','replace')[:300]}")
        text = produced[0].read_text(encoding="utf-8", errors="replace")
    pages = [p.strip() for p in text.split("\f")]
    pages = [p for p in pages if p] or [""]
    return pages, engine, ("explicit-page-break" if len(pages) > 1 else "single-pseudo-page")


# --------------------------------------------------------------------- ZIP
def convert_zip(path: Path) -> tuple[list[str], str, str]:
    """Inventory the pack; extract only documentation members.

    Most of these archives are XBRL taxonomy packages: thousands of .xsd and
    .xml schema files describing regulatory return structures. Dumping their
    bodies into a full-text evidence index produced 705 MB of Markdown and
    100,216 pseudo-pages, against 88,106 real pages in the entire PDF corpus.
    Schema volume would then dominate every topic count, which is exactly the
    frequency bias the analysis is supposed to avoid.

    So: the member listing and a suffix/namespace profile are the evidence.
    Human-readable members (instructions, release notes, guidance) are
    extracted in full. Schema bodies are profiled, not transcribed.
    """
    DOC_SUFFIXES = {".txt", ".md", ".htm", ".html", ".rtf"}
    DOC_NAME_CUES = ("readme", "release", "notes", "instruction", "guidance", "changelog", "licence", "license")
    MAX_DOC_MEMBERS = 40
    try:
        archive = zipfile.ZipFile(path)
    except zipfile.BadZipFile as exc:
        raise RuntimeError(f"BadZipFile: {exc}")
    pages: list[str] = []
    with archive:
        members = [m for m in archive.infolist() if not m.is_dir()]
        by_suffix: dict[str, list] = {}
        for m in members:
            by_suffix.setdefault(Path(m.filename).suffix.lower() or "(none)", []).append(m)

        overview = ["# Archive profile", "",
                    f"Members: {len(members)}",
                    f"Uncompressed bytes: {sum(m.file_size for m in members)}", "",
                    "| suffix | members | bytes |", "|---|---:|---:|"]
        for suffix, group in sorted(by_suffix.items(), key=lambda kv: -len(kv[1])):
            overview.append(f"| {suffix} | {len(group)} | {sum(m.file_size for m in group)} |")

        namespaces: list[str] = []
        for m in [x for x in members if Path(x.filename).suffix.lower() in {".xsd", ".xml"}][:60]:
            try:
                head = archive.read(m)[:4000].decode("utf-8", errors="replace")
            except Exception:
                continue
            for ns in re.findall(r'targetNamespace="([^"]+)"', head):
                if ns not in namespaces:
                    namespaces.append(ns)
        if namespaces:
            overview += ["", "## Target namespaces (sampled)", ""]
            overview += [f"- {ns}" for ns in namespaces[:40]]

        top = sorted(members, key=lambda m: -m.file_size)[:120]
        overview += ["", "## Largest members", "", "| member | bytes |", "|---|---:|"]
        overview += [f"| {m.filename} | {m.file_size} |" for m in top]
        pages.append("\n".join(overview))

        documentation = [
            m for m in members
            if Path(m.filename).suffix.lower() in DOC_SUFFIXES
            or any(cue in Path(m.filename).name.lower() for cue in DOC_NAME_CUES)
        ]
        for m in sorted(documentation, key=lambda x: x.filename.lower())[:MAX_DOC_MEMBERS]:
            if m.file_size > 4_000_000:
                continue
            try:
                body = archive.read(m).decode("utf-8", errors="replace")
            except Exception:
                continue
            if len(body) > MAX_MEMBER_CHARS:
                body = body[:MAX_MEMBER_CHARS] + "\n\n<!-- truncated -->"
            pages.append(f"# {m.filename}\n\n{body}")
        if len(documentation) > MAX_DOC_MEMBERS:
            pages.append(f"<!-- {len(documentation) - MAX_DOC_MEMBERS} further documentation members not transcribed -->")
    return pages, "zipfile-profile", "archive-member"


def convert_docx_resilient(path: Path) -> tuple[list[str], str, str, dict]:
    """OOXML body-order extraction first; LibreOffice only if the package fails."""
    try:
        return convert_docx(path)
    except Exception as first_error:
        try:
            pages, fallback_engine, basis = convert_doc(path)
            return (
                pages,
                f"{fallback_engine} fallback (ooxml: {type(first_error).__name__})",
                basis,
                {"source_text_characters": "", "output_expansion_ratio": ""},
            )
        except Exception as second_error:
            raise RuntimeError(f"ooxml: {first_error}; libreoffice: {second_error}") from second_error


def convert_pdf(path: Path) -> tuple[list[str], str, str]:
    """Real page-anchored extraction, for files whose extension lied."""
    import pymupdf4llm
    chunks = pymupdf4llm.to_markdown(str(path), page_chunks=True, show_progress=False)
    pages = [str(c.get("text") or "") for c in chunks]
    return (pages or [""]), f"pymupdf4llm {getattr(pymupdf4llm, '__version__', '')}".strip(), "source-page"


def sniff(path: Path) -> str:
    """Trust the bytes, not the URL extension.

    The archive mislabels in both directions: four Office files are served with
    a .pdf suffix, and two PDFs are served with a .docx suffix. A converter
    chosen from the extension fails on both.
    """
    head = path.read_bytes()[:8]
    if head[:4] == b"%PDF":
        return "PDF"
    if head[:4] == b"\xd0\xcf\x11\xe0":
        return "DOC"
    if head[:2] == b"PK":
        try:
            with zipfile.ZipFile(path) as z:
                names = set(z.namelist())
            if "word/document.xml" in names:
                return "DOCX"
            if any(n.startswith("ppt/slides/") for n in names):
                return "PPTX"
            if any(n.startswith("xl/") for n in names):
                return "XLSX"
        except zipfile.BadZipFile:
            return "UNKNOWN"
        return "ZIP"
    return "UNKNOWN"


CONVERTERS = {"PDF": convert_pdf, "DOCX": convert_docx_resilient, "PPTX": convert_pptx, "DOC": convert_doc, "ZIP": convert_zip}


def metrics(pages: list[str]) -> dict:
    text = "\n".join(pages)
    nonspace = len(re.sub(r"\s+", "", text))
    normalized_lines = [" ".join(line.split()).casefold() for line in text.splitlines()]
    # Form templates legitimately contain hundreds of repeated empty table rows,
    # underscore writing lines and short labels. They are layout, not duplicated
    # prose. Count only lines with enough semantic content to diagnose the XML
    # alias-expansion failure this metric exists to catch.
    normalized_lines = [
        line for line in normalized_lines
        if sum(character.isalnum() for character in line) >= 20
    ]
    seen: set[str] = set()
    repeated_characters = 0
    for line in normalized_lines:
        if line in seen:
            repeated_characters += len(line)
        else:
            seen.add(line)
    return {"characters": len(text), "nonspace_characters": nonspace,
            "lines": len(text.splitlines()),
            "replacement_characters": text.count("�"),
            "empty_pages": sum(len(re.sub(r"\s+", "", p)) < 30 for p in pages),
            "low_text": nonspace < 100,
            "max_page_characters": max((len(page) for page in pages), default=0),
            "repeated_line_share": round(repeated_characters / max(len(text), 1), 4)}


def render(record: dict, pages: list[str], engine: str, basis: str, m: dict) -> str:
    front = {
        "document_id": f"cbi:{record['sha256']}", "source_url": record["url"],
        "source_urls": record["source_urls"], "source_alias_count": len(record["source_urls"]),
        "source_file": record["localPath"], "source_sha256": record["sha256"],
        "source_bytes": int(record.get("downloadedBytes") or 0),
        "source_pages": len(pages), "extraction_engine": engine,
        "extraction_engine_version": engine, "extraction_pipeline_version": PIPELINE_VERSION,
        "detected_source_format": record["format"], "page_basis": basis,
        "ocr_enabled": False, "quality_low_text": bool(m["low_text"]),
        "quality_empty_pages": int(m["empty_pages"]),
    }
    out = ["---"] + [f"{k}: {yaml_scalar(v)}" for k, v in front.items()] + ["---", ""]
    for n, page in enumerate(pages, 1):
        out += [f"<!-- source-page: {n} -->", f'<a id="page-{n}"></a>', "", page.strip(), ""]
    return "\n".join(out).rstrip() + "\n"


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--archive", type=Path, required=True)
    ap.add_argument("--output", type=Path, required=True)
    ap.add_argument("--formats", default="DOCX,DOC,PPTX,ZIP")
    ap.add_argument("--max-files", type=int, default=0)
    ap.add_argument("--force", action="store_true",
                    help="reprocess every selected source even when the journal records success")
    ap.add_argument("--force-sha", action="append", default=[], metavar="SHA256",
                    help="reprocess this source hash even if the journal records success; repeatable")
    args = ap.parse_args()
    archive, root = args.archive.resolve(), args.output.resolve()
    wanted = {f.strip().upper() for f in args.formats.split(",")}
    forced = {value.strip().lower() for value in args.force_sha}

    with (archive / "manifests" / "files.csv").open(encoding="utf-8-sig", newline="") as fh:
        records = [r for r in csv.DictReader(fh)
                   if r["status"] == "downloaded" and r["format"].upper() in wanted]
    grouped: dict[str, dict] = {}
    for r in sorted(records, key=lambda r: r["url"]):
        g = grouped.setdefault(r["sha256"], {**r, "source_urls": []})
        g["source_urls"].append(r["url"])
    todo = list(grouped.values())
    if args.max_files:
        todo = todo[: args.max_files]
    print(f"logical office/archive files: {len(todo)}", flush=True)

    journal_path = root / "conversion-journal.jsonl"
    done: dict[str, dict] = {}
    if journal_path.exists():
        for line in journal_path.read_text(encoding="utf-8").splitlines():
            if line.strip():
                row = json.loads(line)
                if row.get("status") in {"success", "low_text"}:
                    done[row["source_sha256"]] = row
    root.mkdir(parents=True, exist_ok=True)
    started = time.time()
    with journal_path.open("a", encoding="utf-8") as journal:
        for i, record in enumerate(todo, 1):
            if record["sha256"] in done and not args.force and record["sha256"].lower() not in forced:
                continue
            previous_success = done.get(record["sha256"])
            source = archive / record["localPath"].replace("\\", "/")
            destination = root / "markdown" / record["sha256"][:2] / record["sha256"][2:4] / f"{record['sha256']}.md"
            row = {"url": record["url"], "source_urls": record["source_urls"],
                   "alias_count": len(record["source_urls"]), "source_file": record["localPath"],
                   "source_sha256": record["sha256"], "source_bytes": record.get("downloadedBytes"),
                   "format": record["format"].upper(),
                   "pipeline_version": PIPELINE_VERSION,
                   "markdown_file": str(destination.relative_to(root)).replace("\\", "/"),
                   "status": "success", "error": None}
            try:
                if not source.is_file():
                    raise FileNotFoundError(source)
                declared = record["format"].upper()
                detected = sniff(source)
                effective = detected if detected in CONVERTERS else declared
                if effective != declared:
                    row["format_mismatch"] = f"declared={declared} detected={detected}"
                    record = {**record, "format": f"{detected} (served as {declared})"}
                diagnostics = {}
                converted = CONVERTERS[effective](source)
                if effective == "DOCX":
                    pages, engine, basis, diagnostics = converted
                else:
                    pages, engine, basis = converted
                m = metrics(pages)
                m.update(diagnostics)
                atomic_text_write(destination, render(record, pages, engine, basis, m))
                row.update(m)
                row.update({"engine": engine, "page_basis": basis, "pages": len(pages),
                            "markdown_bytes": destination.stat().st_size,
                            "markdown_sha256": hashlib.sha256(destination.read_bytes()).hexdigest(),
                            "status": "low_text" if m["low_text"] else "success"})
            except Exception as exc:
                row.update({"status": "error", "error": f"{type(exc).__name__}: {exc}"[:800],
                            "engine": "", "page_basis": "", "pages": 0})
            journal.write(json.dumps(row, ensure_ascii=False) + "\n"); journal.flush()
            # A failed forced refresh must not erase a previously successful,
            # still-valid conversion from the materialised manifest. The failed
            # attempt remains in the append-only journal for diagnosis.
            if row["status"] in {"success", "low_text"} or previous_success is None:
                done[record["sha256"]] = row
            if i % 10 == 0 or i == len(todo):
                print(f"  {i}/{len(todo)}  {time.time()-started:.0f}s  last={row['status']}", flush=True)

    # The journal is the durable state. A targeted ``--formats`` or
    # ``--force-sha`` run must not truncate an existing all-format manifest to
    # only the current selection.
    rows = list(done.values())
    fields = ["url", "source_urls", "alias_count", "source_file", "source_sha256", "source_bytes",
              "format", "pipeline_version", "markdown_file", "markdown_bytes", "markdown_sha256", "engine", "page_basis",
              "pages", "characters", "nonspace_characters", "lines", "replacement_characters",
              "empty_pages", "low_text", "max_page_characters", "repeated_line_share",
              "source_text_characters", "output_expansion_ratio",
              "status", "error", "format_mismatch"]
    with (root / "conversion-manifest.csv").open("w", encoding="utf-8", newline="") as fh:
        w = csv.DictWriter(
            fh, fieldnames=fields, extrasaction="ignore", lineterminator="\n")
        w.writeheader()
        for r in sorted(rows, key=lambda r: r["url"]):
            r = dict(r); r["source_urls"] = " | ".join(r.get("source_urls") or [])
            w.writerow(r)
    summary = {"generated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
               "pipeline_version": PIPELINE_VERSION, "records": len(rows),
               "by_format": {}, "statuses": {}, "by_page_basis": {},
               "pages": sum(int(r.get("pages") or 0) for r in rows),
               "source_bytes": sum(int(r.get("source_bytes") or 0) for r in rows),
               "markdown_bytes": sum(int(r.get("markdown_bytes") or 0) for r in rows)}
    for r in rows:
        for key, field in (("by_format", "format"), ("statuses", "status"), ("by_page_basis", "page_basis")):
            summary[key][r.get(field) or "n/a"] = summary[key].get(r.get(field) or "n/a", 0) + 1
    (root / "conversion-summary.json").write_text(
        json.dumps(summary, indent=2) + "\n", encoding="utf-8", newline="\n")
    print(json.dumps(summary, indent=2), flush=True)
    return 0


if __name__ == "__main__":
    sys.exit(main())
